from __future__ import annotations

import json
from pathlib import Path

import pytest

from uagent.tools.office_to_markdown_tool import run_tool


def load_result(out: str) -> dict:
    result = json.loads(out)
    assert isinstance(result, dict)
    return result


def test_rejects_unsupported_file(repo_tmp_path: Path) -> None:
    path = repo_tmp_path / "sample.txt"
    path.write_text("hello", encoding="utf-8")
    result = load_result(run_tool({"input_path": str(path)}))
    assert result["ok"] is False


def test_converts_docx(repo_tmp_path: Path) -> None:
    pytest.importorskip("docx")
    from docx import Document

    path = repo_tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("Title", level=1)
    doc.add_paragraph("Body")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    doc.save(path)

    result = load_result(run_tool({"input_path": str(path)}))
    assert result["ok"] is True
    assert "# Title" in result["markdown"]
    assert "| A | B |" in result["markdown"]


def test_converts_xlsx(repo_tmp_path: Path) -> None:
    pytest.importorskip("exstruct")
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    path = repo_tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Value"])
    ws.append(["A", 1])
    wb.save(path)
    wb.close()

    output = repo_tmp_path / "converted.md"
    result = load_result(
        run_tool({"input_path": str(path), "output_path": str(output)})
    )
    assert result["ok"] is True
    assert output.read_text(encoding="utf-8").startswith("# sample")
    assert "## Data" in output.read_text(encoding="utf-8")


def test_converts_pptx(repo_tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    path = repo_tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Slide title"
    presentation.save(path)

    result = load_result(run_tool({"input_path": str(path), "include_notes": False}))
    assert result["ok"] is True
    assert "## Slide 1" in result["markdown"]
    assert "Slide title" in result["markdown"]
