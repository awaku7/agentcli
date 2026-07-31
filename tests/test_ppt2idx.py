import sys
from pathlib import Path
import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))

import pptx
from uagent.tools.ppt2idx_tool import run_tool


@pytest.fixture
def sample_pptx():
    tmp_dir = Path("tmp_test_dir")
    tmp_dir.mkdir(exist_ok=True)

    prs = pptx.Presentation()

    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Sample Title"
    slide1.placeholders[1].text = "Sample Subtitle"
    slide1.notes_slide.notes_text_frame.text = "Sample Note"

    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Second Slide"
    slide2.placeholders[1].text = "Content line 1\nContent line 2"

    path = tmp_dir / "test.pptx"
    prs.save(str(path))
    yield str(path)

    if path.exists():
        path.unlink()
    if tmp_dir.exists():
        try:
            tmp_dir.rmdir()
        except Exception:
            pass


def test_ppt2idx_index(sample_pptx):
    res = run_tool({"path": sample_pptx, "mode": "index"})
    assert ("Index for:" in res) or ("インデックス:" in res)
    assert "Slide  1: Sample Title - Sample Subtitle [Notes]" in res
    assert "Slide  2: Second Slide - Content line 1 Content line 2" in res


def test_ppt2idx_section(sample_pptx):
    res1 = run_tool({"path": sample_pptx, "mode": "section", "section": 1})
    assert "=== Slide 1: Sample Title ===" in res1
    assert "Sample Subtitle" in res1
    assert "-- Speaker Notes --" in res1
    assert "Sample Note" in res1

    res2 = run_tool({"path": sample_pptx, "mode": "section", "section": 2})
    assert "=== Slide 2: Second Slide ===" in res2
    assert "Content line 1" in res2


def test_ppt2idx_errors():
    res_err_path = run_tool({"mode": "index"})
    assert ("Error: 'path' is required" in res_err_path) or (
        "エラー: 'path' は必須です。" in res_err_path
    )

    res_err_file = run_tool({"path": "non_existent.pptx", "mode": "index"})
    assert ("Error: File not found" in res_err_file) or (
        "エラー: ファイルが見つかりません" in res_err_file
    )

    res_err_sec = run_tool({"path": "non_existent.pptx", "mode": "section"})
    assert ("Error: File not found" in res_err_sec) or (
        "エラー: ファイルが見つかりません" in res_err_sec
    )
