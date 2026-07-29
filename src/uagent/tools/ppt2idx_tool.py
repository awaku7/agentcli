from __future__ import annotations

import os
from typing import Any, List, Dict

from .._pip_auto import install_with_status

if install_with_status("python-pptx", "pptx", version_spec=">=1.0.2"):
    import pptx
else:
    pptx = None

from .i18n_helper import make_tool_translator
from .index_tool_helpers import resolve_index_path

_ = make_tool_translator(__file__)

TOOL_SPEC = {
    "type": "function",
    "tool_genre": "index",
    "function": {
        "name": "ppt2idx",
        "description": _(
            "tool.description",
            default=(
                "Parse a PowerPoint presentation (.pptx) into slide titles and text summaries, "
                "returning a numbered index or a specific slide section. Use this when reading large PowerPoint files: "
                "call with mode='index' to get the list of slides with titles and preview text, then call with "
                "mode='section' and section=N to retrieve the full text content of that slide."
            ),
        ),
        "x_search_terms": _(
            "x_search_terms",
            default=[
                "read pptx file",
                "powerpoint index",
                "powerpoint parser",
                "powerpoint slide reader",
                "ppt section",
                "PPTXファイルを読む",
                "PowerPointインデックス",
                "スライド目次",
                "スライド抽出",
                "セクション分割",
            ],
        ),
        "x_search_terms_en": [
            "read pptx file",
            "powerpoint index",
            "powerpoint parser",
            "powerpoint slide reader",
            "ppt section",
        ],
        "x_parallel_safe": True,
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": _(
                        "param.path.description",
                        default="Path to the PowerPoint (.pptx) file.",
                    ),
                },
                "mode": {
                    "type": "string",
                    "enum": ["index", "section"],
                    "description": _(
                        "param.mode.description",
                        default=(
                            '"index" returns a numbered table of contents listing slides with titles and summaries. '
                            '"section" returns the full text content of a specific slide.'
                        ),
                    ),
                },
                "section": {
                    "type": "integer",
                    "description": _(
                        "param.section.description",
                        default=(
                            "Slide number to retrieve (used only when mode='section'). "
                            "Get the number from the index output."
                        ),
                    ),
                },
            },
            "required": ["path", "mode"],
            "additionalProperties": False,
        },
    },
}


class _PptxIndexBuilder:
    def __init__(self, filepath: str):
        self.filepath = filepath
        self.slides: List[Dict[str, Any]] = []
        self._parse()

    def _parse(self) -> None:
        if pptx is None:
            raise RuntimeError("python-pptx library is required to parse .pptx files.")

        prs = pptx.Presentation(self.filepath)

        for idx, slide in enumerate(prs.slides, 1):
            title = ""
            texts: List[str] = []
            notes = ""

            # Check notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    notes = notes_text

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                full_text = tf.text.strip()
                if not full_text:
                    continue

                if shape == slide.shapes.title or (
                    hasattr(shape, "is_placeholder")
                    and shape.is_placeholder
                    and shape.placeholder_format.idx == 0
                ):
                    if not title:
                        title = full_text.replace("\n", " ")
                texts.append(full_text)

            if not title and texts:
                title = texts[0].replace("\n", " ")
            if not title:
                title = "(Untitled Slide)"

            # Build preview text
            combined_body = " / ".join(
                [t.replace("\n", " ") for t in texts if t.replace("\n", " ") != title]
            )
            preview = combined_body[:60] + ("..." if len(combined_body) > 60 else "")

            self.slides.append(
                {
                    "slide_num": idx,
                    "title": title,
                    "preview": preview,
                    "texts": texts,
                    "notes": notes,
                }
            )


def run_tool(args: dict[str, Any]) -> str:
    path = args.get("path")
    if not path:
        return _("err.path_required", default="Error: 'path' is required.")

    mode = args.get("mode", "index")
    section = args.get("section")

    try:
        resolved = resolve_index_path(path)
    except Exception as e:
        return _("err.read_error", default="Error reading file: {e}").format(e=e)

    if not os.path.isfile(resolved):
        return _("err.file_not_found", default="Error: File not found: {path}").format(
            path=path
        )

    try:
        builder = _PptxIndexBuilder(resolved)
    except Exception as e:
        return _(
            "err.parse_error", default="Error parsing PowerPoint file: {e}"
        ).format(e=e)

    if not builder.slides:
        return _("msg.no_slides", default="(no slides found in presentation)")

    if mode == "index":
        toc_lines = []
        for slide in builder.slides:
            s_num = slide["slide_num"]
            title = slide["title"]
            prev = slide["preview"]
            prev_str = f" - {prev}" if prev else ""
            notes_flag = " [Notes]" if slide["notes"] else ""
            toc_lines.append(f"Slide {s_num:2d}: {title}{prev_str}{notes_flag}")

        toc = "\n".join(toc_lines)
        return _(
            "msg.index_output",
            default=(
                "Index for: {path}\n---\n{toc}\n---\n"
                "Total slides: {total}\n"
                "To retrieve a slide, call ppt2idx with mode='section' and section=N."
            ),
        ).format(path=path, toc=toc, total=len(builder.slides))

    elif mode == "section":
        if section is None:
            return _(
                "err.section_required",
                default="Error: 'section' (integer) is required when mode='section'.",
            )
        try:
            section_num = int(section)
        except (ValueError, TypeError):
            return _(
                "err.section_invalid", default="Error: 'section' must be an integer."
            )

        if section_num < 1 or section_num > len(builder.slides):
            return _(
                "err.section_not_found",
                default="Error: Slide {section_num} not found. Valid range: 1..{last}.",
            ).format(section_num=section_num, last=len(builder.slides))

        target = builder.slides[section_num - 1]
        out_lines = [f"=== Slide {target['slide_num']}: {target['title']} ==="]
        if target["texts"]:
            out_lines.append("\n-- Content --")
            out_lines.extend(target["texts"])
        if target["notes"]:
            out_lines.append("\n-- Speaker Notes --")
            out_lines.append(target["notes"])

        return "\n".join(out_lines)

    else:
        return _(
            "err.invalid_mode",
            default="Error: Invalid mode '{mode}'. Use 'index' or 'section'.",
        ).format(mode=mode)
