# tools/read_pptx_pdf.py
# -*- coding: utf-8 -*-

"""Safe reading of PDF and PPTX files.

This module provides compatibility layers and extractors for PDF and PPTX files,
mapping them to a common JSON schema.
"""

# --- imports at top ---
import collections
import collections.abc
import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import logging

import unicodedata
from collections import Counter

from .i18n_helper import make_tool_translator

_ = make_tool_translator(__file__)

# Optional external libraries
try:
    import pdfplumber
except ImportError:
    pdfplumber = None  # type: ignore[assignment]

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None  # type: ignore[assignment]
    MSO_SHAPE_TYPE = None  # type: ignore[assignment]

# Compatibility layer: complement collections and collections.abc
for name in ("Mapping", "MutableMapping", "Sequence"):
    if not hasattr(collections, name) and hasattr(collections.abc, name):
        setattr(collections, name, getattr(collections.abc, name))

# Suppress PDF related logs
logging.getLogger("pdfminer").setLevel(logging.ERROR)

BUSY_LABEL = True
STATUS_LABEL = "tool:read_pptx_pdf"

DEFAULT_MAX_CHARS = 8000
JSON_SCHEMA_VERSION = "1.1"


TOOL_SPEC: Dict[str, Any] = {
    "type": "function",
    "function": {
        "name": "read_pptx_pdf",
        "description": _(
            "tool.description",
            default=(
                "Read PDF/PPTX (or their common JSON schema) and return extracted text by page. "
                "You can set path to .pdf / .pptx / .json. For PDF/PPTX, this module converts it into the common JSON schema."
            ),
        ),
        "system_prompt": _(
            "tool.system_prompt",
            default=(
                "IMPORTANT: This tool reads PDF/PPTX/JSON and returns page-level text.\n"
                "Input: path, page_index, max_chars\n"
                "Output: page text (all pages or specified page)\n\n"
            ),
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": _(
                        "param.path.description",
                        default="Path to a PDF / PPTX / JSON file.",
                    ),
                },
                "page_index": {
                    "type": "integer",
                    "description": _(
                        "param.page_index.description",
                        default=(
                            "1-based page number (PDF page or PPTX slide index). "
                            "If omitted, returns concatenated text for all pages."
                        ),
                    ),
                },
                "max_chars": {
                    "type": "integer",
                    "description": _(
                        "param.max_chars.description",
                        default=(
                            "Maximum number of characters to return. If omitted, truncates at 8000 characters."
                        ),
                    ),
                },
            },
            "required": ["path"],
        },
    },
}


# ==============================
# Utilities
# ==============================


def normalize_text(s: str) -> str:
    """Normalize text extracted from PDF / PPTX."""
    if not s:
        return ""
    # Remove invisible characters
    for ch in ("\ufeff", "\u200b", "\u200c", "\u200d"):
        s = s.replace(ch, "")
    # NFKC normalization
    s = unicodedata.normalize("NFKC", s)
    return s


def infer_style_from_fontnames(
    fontnames: List[Optional[str]],
) -> Optional[Dict[str, Optional[bool]]]:
    """Roughly estimate bold/italic styles from font names."""
    names = [fn.lower() for fn in fontnames if fn]
    if not names:
        return None

    bold_hits = [
        n for n in names if any(k in n for k in ("bold", "black", "heavy", "demi"))
    ]
    italic_hits = [
        n for n in names if any(k in n for k in ("italic", "oblique", "ita", "kursiv"))
    ]

    style: Dict[str, Optional[bool]] = {
        "bold": bool(bold_hits),
        "italic": bool(italic_hits),
    }
    return style


def summarize_font_from_words(
    words: List[Dict[str, Any]],
) -> Optional[Dict[str, Optional[Any]]]:
    """Determine representative font name and size from pdfplumber word dicts."""
    if not words:
        return None

    name_counter: Counter = Counter()
    size_counter: Counter = Counter()

    for w in words:
        fn = w.get("fontname")
        if fn:
            name_counter[fn] += 1
        sz = w.get("size")
        if sz:
            try:
                size_counter[round(float(sz), 1)] += 1
            except Exception:
                pass

    font_name: Optional[str] = None
    font_size: Optional[float] = None

    if name_counter:
        font_name = name_counter.most_common(1)[0][0]
    if size_counter:
        font_size = size_counter.most_common(1)[0][0]

    if font_name is None and font_size is None:
        return None

    return {"name": font_name, "size": font_size}


def estimate_pdf_alignment(
    left: float, right: float, page_width: float
) -> Optional[str]:
    """Roughly estimate alignment (left / center / right) from bounding box."""
    if page_width <= 0:
        return None

    center_x = (left + right) / 2.0
    page_center = page_width / 2.0

    if abs(center_x - page_center) <= page_width * 0.05:
        return "center"

    margin_left = left
    margin_right = page_width - right

    if margin_right < margin_left * 0.5:
        return "right"

    return "left"


# ==============================
# PDF -> blocks
# ==============================


def build_pdf_blocks(
    words_raw: List[Dict[str, Any]], page_width: float
) -> List[Dict[str, Any]]:
    """Build paragraph blocks from pdfplumber.extract_words() results."""

    norm_words: List[Dict[str, Any]] = []
    for w in words_raw:
        t = normalize_text(str(w.get("text", "")))
        if not t.strip():
            continue
        w2 = dict(w)
        w2["text"] = t
        norm_words.append(w2)

    if not norm_words:
        return []

    norm_words.sort(key=lambda x: (x.get("top", 0), x.get("x0", 0)))

    lines: List[List[Dict[str, Any]]] = []
    current: List[Dict[str, Any]] = []
    current_top: Optional[float] = None

    def _same_line(a: float, b: float) -> bool:
        return abs(a - b) <= 2.0

    for w in norm_words:
        top = float(w.get("top", 0.0) or 0.0)
        if current_top is None:
            current_top = top
            current = [w]
            continue

        if _same_line(top, current_top):
            current.append(w)
        else:
            lines.append(current)
            current = [w]
            current_top = top

    if current:
        lines.append(current)

    line_objs: List[Dict[str, Any]] = []
    for ln in lines:
        ln.sort(key=lambda x: float(x.get("x0", 0.0) or 0.0))
        text = " ".join([str(w.get("text", "")) for w in ln]).strip()
        if not text:
            continue
        left = float(min([w.get("x0", 0.0) or 0.0 for w in ln]))
        right = float(max([w.get("x1", 0.0) or 0.0 for w in ln]))
        top = float(min([w.get("top", 0.0) or 0.0 for w in ln]))
        bottom = float(max([w.get("bottom", 0.0) or 0.0 for w in ln]))
        line_objs.append(
            {
                "text": text,
                "bbox": [left, top, right, bottom],
                "words": ln,
            }
        )

    if not line_objs:
        return []

    blocks: List[Dict[str, Any]] = []
    cur_block: List[Dict[str, Any]] = [line_objs[0]]

    def _gap(a: Dict[str, Any], b: Dict[str, Any]) -> float:
        return float(b["bbox"][1]) - float(a["bbox"][3])

    for prev, nxt in zip(line_objs, line_objs[1:]):
        g = _gap(prev, nxt)
        if g > 8.0:
            blocks.append({"lines": cur_block})
            cur_block = [nxt]
        else:
            cur_block.append(nxt)

    if cur_block:
        blocks.append({"lines": cur_block})

    for b in blocks:
        b_lines = b.get("lines") or []
        if not b_lines:
            continue

        left = min([float(ln["bbox"][0]) for ln in b_lines])
        right = max([float(ln["bbox"][2]) for ln in b_lines])
        top = min([float(ln["bbox"][1]) for ln in b_lines])
        bottom = max([float(ln["bbox"][3]) for ln in b_lines])
        b["bbox"] = [left, top, right, bottom]
        b["align"] = estimate_pdf_alignment(left, right, page_width)

        words_all: List[Dict[str, Any]] = []
        for ln in b_lines:
            words_all.extend(list(ln.get("words") or []))
        b["font"] = summarize_font_from_words(words_all)

        fontnames = [w.get("fontname") for w in words_all]
        b["style"] = infer_style_from_fontnames(fontnames)

        b["text"] = "\n".join([str(ln.get("text", "")) for ln in b_lines]).strip()

    return blocks


# ==============================
# JSON schema helpers
# ==============================


def _ensure_common_schema(doc: Dict[str, Any]) -> Dict[str, Any]:
    """Ensure minimal common schema fields exist."""
    if "schema_version" not in doc:
        doc["schema_version"] = JSON_SCHEMA_VERSION
    if "pages" not in doc or not isinstance(doc["pages"], list):
        doc["pages"] = []
    return doc


def _wrap_page_texts(pages: List[str]) -> Dict[str, Any]:
    return _ensure_common_schema(
        {
            "schema_version": JSON_SCHEMA_VERSION,
            "pages": [{"index": i + 1, "text": t} for i, t in enumerate(pages)],
        }
    )


# ==============================
# Extractors
# ==============================


def _extract_pdf_pages(pdf_path: str) -> Tuple[List[str], List[str]]:
    warnings: List[str] = []

    if pdfplumber is None:
        warnings.append("pdfplumber is not available")
        return [], warnings

    pages_text: List[str] = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                try:
                    words = page.extract_words(
                        use_text_flow=True,
                        keep_blank_chars=False,
                        extra_attrs=["fontname", "size"],
                    )
                except Exception:
                    words = []

                try:
                    blocks = build_pdf_blocks(words, page.width)
                    txt = "\n\n".join(
                        [b.get("text", "") for b in blocks if b.get("text")]
                    )
                    txt = normalize_text(txt)
                    if not txt.strip():
                        txt = normalize_text(page.extract_text() or "")
                except Exception:
                    txt = normalize_text(page.extract_text() or "")

                pages_text.append(txt)
    except Exception as e:
        warnings.append(f"pdfplumber failed: {type(e).__name__}: {e}")
        return [], warnings

    return pages_text, warnings


def _extract_pptx_pages(pptx_path: str) -> Tuple[List[str], List[str]]:
    warnings: List[str] = []

    if Presentation is None:
        warnings.append("python-pptx is not available")
        return [], warnings

    pages_text: List[str] = []
    try:
        pres = Presentation(pptx_path)
        for slide in pres.slides:
            parts: List[str] = []
            for shape in slide.shapes:
                try:
                    if (
                        MSO_SHAPE_TYPE is not None
                        and shape.shape_type == MSO_SHAPE_TYPE.GROUP
                    ):
                        continue
                except Exception:
                    pass

                txt = ""
                try:
                    if (
                        getattr(shape, "has_text_frame", False)
                        and shape.text_frame is not None
                    ):
                        txt = shape.text_frame.text or ""
                except Exception:
                    txt = ""

                txt = normalize_text(txt)
                if txt.strip():
                    parts.append(txt)

            pages_text.append("\n".join(parts).strip())
    except Exception as e:
        warnings.append(f"python-pptx failed: {type(e).__name__}: {e}")
        return [], warnings

    return pages_text, warnings


# ==============================
# Main
# ==============================


def _read_common_json(path: str) -> Tuple[Dict[str, Any], List[str]]:
    warnings: List[str] = []
    try:
        data = json.loads(Path(path).read_text(encoding="utf-8"))
        if not isinstance(data, dict):
            return _ensure_common_schema({}), ["JSON root is not an object"]
        return _ensure_common_schema(data), warnings
    except Exception as e:
        return _ensure_common_schema({}), [
            f"Failed to read JSON: {type(e).__name__}: {e}"
        ]


def _get_pages_text(doc: Dict[str, Any]) -> List[str]:
    pages = doc.get("pages")
    if not isinstance(pages, list):
        return []
    out: List[str] = []
    for p in pages:
        if not isinstance(p, dict):
            continue
        t = p.get("text")
        out.append(normalize_text(str(t or "")))
    return out


def run_tool(args: Dict[str, Any]) -> str:
    path = (args.get("path") or "").strip()
    page_index = args.get("page_index")
    max_chars = args.get("max_chars")

    if not path:
        return "[read_pptx_pdf error] path is required"

    try:
        if max_chars is None:
            max_chars_i = DEFAULT_MAX_CHARS
        else:
            max_chars_i = int(max_chars)
            if max_chars_i <= 0:
                max_chars_i = DEFAULT_MAX_CHARS
    except Exception:
        max_chars_i = DEFAULT_MAX_CHARS

    p = Path(path)
    if not p.exists() or not p.is_file():
        return f"[read_pptx_pdf error] file not found: {path}"

    suffix = p.suffix.lower()
    warnings: List[str] = []

    if suffix == ".json":
        doc, warnings = _read_common_json(path)
        pages_text = _get_pages_text(doc)
    elif suffix == ".pdf":
        pages_text, warnings = _extract_pdf_pages(path)
        doc = _wrap_page_texts(pages_text)
    elif suffix == ".pptx":
        pages_text, warnings = _extract_pptx_pages(path)
        doc = _wrap_page_texts(pages_text)
    else:
        return f"[read_pptx_pdf error] unsupported file extension: {suffix}"

    pages_text = _get_pages_text(doc)

    if page_index is not None:
        try:
            idx = int(page_index)
        except Exception:
            return "[read_pptx_pdf error] page_index must be an integer"

        if idx <= 0 or idx > len(pages_text):
            return f"[read_pptx_pdf error] page_index out of range: {idx}"

        out = pages_text[idx - 1]
    else:
        out = "\n\n".join([t for t in pages_text if t])

    if not out.strip():
        warnings.append("No text could be extracted")

    out = out[:max_chars_i]

    if warnings:
        warn_text = "\n".join([f"[read_pptx_pdf warn] {w}" for w in warnings])
        if out:
            return warn_text + "\n" + out
        return warn_text

    return out
