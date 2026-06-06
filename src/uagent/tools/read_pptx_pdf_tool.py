from __future__ import annotations

# tools/read_pptx_pdf.py
# -*- coding: utf-8 -*-

"""Safe reading of PDF and PPTX files.

This module provides compatibility layers and extractors for PDF and PPTX files,
mapping them to a common JSON schema.
"""

# --- imports at top ---
import collections
import collections.abc
import json
from io import BytesIO
from pathlib import Path
from typing import Any, Optional

import logging

import unicodedata
from collections import Counter

from .i18n_helper import make_tool_translator

_ = make_tool_translator(__file__)

# Optional external libraries
try:
    import pdfplumber
except ImportError:
    pdfplumber = None  # type: ignore[assignment]

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None  # type: ignore[assignment]
    MSO_SHAPE_TYPE = None  # type: ignore[assignment]

try:
    import msoffcrypto
except Exception:  # pragma: no cover
    msoffcrypto = None  # type: ignore[assignment]

# Compatibility layer: complement collections and collections.abc
for name in ("Mapping", "MutableMapping", "Sequence"):
    if not hasattr(collections, name) and hasattr(collections.abc, name):
        setattr(collections, name, getattr(collections.abc, name))

# Suppress PDF related logs
logging.getLogger("pdfminer").setLevel(logging.ERROR)

BUSY_LABEL = True
STATUS_LABEL = "tool:read_pptx_pdf"

DEFAULT_MAX_CHARS = 8000
JSON_SCHEMA_VERSION = "1.1"


TOOL_SPEC: dict[str, Any] = {
    "tool_level": 1,
    "tool_genre": "office",
    "type": "function",
    "function": {
        "name": "read_pptx_pdf",
        "description": _(
            "tool.description",
            default=(
                "Read PDF/PPTX (or their common JSON schema) and return extracted text by page. "
                "You can set path to .pdf / .pptx / .json. For PDF/PPTX, this module converts it into the common JSON schema."
            ),
        ),
        "system_prompt": _(
            "tool.system_prompt",
            default=(
                "IMPORTANT: This tool reads PDF/PPTX/JSON and returns page-level text.\n"
                "Input: path, page_index, max_chars\n"
                "Output: page text (all pages or specified page)\n\n"
            ),
        ),
        "x_search_terms": _(
            "x_search_terms",
            default=[
                "read_pptx_pdf",
                "read pptx pdf",
                "pdf",
                "pptx",
                "page text",
                "extract pages",
            ],
        ),
        "x_search_terms_en": [
            "read_pptx_pdf",
            "read pptx pdf",
            "pdf",
            "pptx",
            "page text",
            "extract pages",
        ],
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": _(
                        "param.path.description",
                        default="Path to a PDF / PPTX / JSON file.",
                    ),
                },
                "password": {
                    "type": "string",
                    "description": _(
                        "param.password.description",
                        default=(
                            "Optional password for encrypted PDF/PPTX files. If omitted and the input is encrypted, "
                            "the tool will prompt once for a password."
                        ),
                    ),
                },
                "page_index": {
                    "type": "integer",
                    "description": _(
                        "param.page_index.description",
                        default=(
                            "1-based page number (PDF page or PPTX slide index). "
                            "If omitted, returns concatenated text for all pages. Alias of page."
                        ),
                    ),
                },
                "page": {
                    "type": "integer",
                    "description": _(
                        "param.page.description",
                        default="Page number to read (1-based). Alias of page_index.",
                    ),
                },
                "max_chars": {
                    "type": "integer",
                    "description": _(
                        "param.max_chars.description",
                        default=(
                            "Maximum number of characters to return. If omitted, truncates at 8000 characters."
                        ),
                    ),
                },
            },
            "required": ["path"],
        },
    },
}


# ==============================
# Utilities
# ==============================


def normalize_text(s: str) -> str:
    """Normalize text extracted from PDF / PPTX."""
    if not s:
        return ""
    # Remove invisible characters
    for ch in ("\ufeff", "\u200b", "\u200c", "\u200d"):
        s = s.replace(ch, "")
    # NFKC normalization
    s = unicodedata.normalize("NFKC", s)
    return s


def infer_style_from_fontnames(
    fontnames: list[Optional[str]],
) -> Optional[dict[str, Optional[bool]]]:
    """Roughly estimate bold/italic styles from font names."""
    names = [fn.lower() for fn in fontnames if fn]
    if not names:
        return None

    bold_hits = [
        n for n in names if any(k in n for k in ("bold", "black", "heavy", "demi"))
    ]
    italic_hits = [
        n for n in names if any(k in n for k in ("italic", "oblique", "ita", "kursiv"))
    ]

    style: dict[str, Optional[bool]] = {
        "bold": bool(bold_hits),
        "italic": bool(italic_hits),
    }
    return style


def summarize_font_from_words(
    words: list[dict[str, Any]],
) -> Optional[dict[str, Optional[Any]]]:
    """Determine representative font name and size from pdfplumber word dicts."""
    if not words:
        return None

    name_counter: Counter = Counter()
    size_counter: Counter = Counter()

    for w in words:
        fn = w.get("fontname")
        if fn:
            name_counter[fn] += 1
        sz = w.get("size")
        if sz:
            try:
                size_counter[round(float(sz), 1)] += 1
            except Exception:
                pass

    font_name: Optional[str] = None
    font_size: Optional[float] = None

    if name_counter:
        font_name = name_counter.most_common(1)[0][0]
    if size_counter:
        font_size = size_counter.most_common(1)[0][0]

    if font_name is None and font_size is None:
        return None

    return {"name": font_name, "size": font_size}


def estimate_pdf_alignment(
    left: float, right: float, page_width: float
) -> Optional[str]:
    """Roughly estimate alignment (left / center / right) from bounding box."""
    if page_width <= 0:
        return None

    center_x = (left + right) / 2.0
    page_center = page_width / 2.0

    if abs(center_x - page_center) <= page_width * 0.05:
        return "center"

    margin_left = left
    margin_right = page_width - right

    if margin_right < margin_left * 0.5:
        return "right"

    return "left"


# ==============================
# PDF -> blocks
# ==============================


def build_pdf_blocks(
    words_raw: list[dict[str, Any]], page_width: float
) -> list[dict[str, Any]]:
    """Build paragraph blocks from pdfplumber.extract_words() results."""

    norm_words: list[dict[str, Any]] = []
    for w in words_raw:
        t = normalize_text(str(w.get("text", "")))
        if not t.strip():
            continue
        w2 = dict(w)
        w2["text"] = t
        norm_words.append(w2)

    if not norm_words:
        return []

    norm_words.sort(key=lambda x: (x.get("top", 0), x.get("x0", 0)))

    lines: list[list[dict[str, Any]]] = []
    current: list[dict[str, Any]] = []
    current_top: Optional[float] = None

    def _same_line(a: float, b: float) -> bool:
        return abs(a - b) <= 2.0

    for w in norm_words:
        top = float(w.get("top", 0.0) or 0.0)
        if current_top is None:
            current_top = top
            current = [w]
            continue

        if _same_line(top, current_top):
            current.append(w)
        else:
            lines.append(current)
            current = [w]
            current_top = top

    if current:
        lines.append(current)

    line_objs: list[dict[str, Any]] = []
    for ln in lines:
        ln.sort(key=lambda x: float(x.get("x0", 0.0) or 0.0))
        text = " ".join([str(w.get("text", "")) for w in ln]).strip()
        if not text:
            continue
        left = float(min([w.get("x0", 0.0) or 0.0 for w in ln]))
        right = float(max([w.get("x1", 0.0) or 0.0 for w in ln]))
        top = float(min([w.get("top", 0.0) or 0.0 for w in ln]))
        bottom = float(max([w.get("bottom", 0.0) or 0.0 for w in ln]))
        line_objs.append(
            {
                "text": text,
                "bbox": [left, top, right, bottom],
                "words": ln,
            }
        )

    if not line_objs:
        return []

    blocks: list[dict[str, Any]] = []
    cur_block: list[dict[str, Any]] = [line_objs[0]]

    def _gap(a: dict[str, Any], b: dict[str, Any]) -> float:
        return float(b["bbox"][1]) - float(a["bbox"][3])

    for prev, nxt in zip(line_objs, line_objs[1:], strict=False):
        g = _gap(prev, nxt)
        if g > 8.0:
            blocks.append({"lines": cur_block})
            cur_block = [nxt]
        else:
            cur_block.append(nxt)

    if cur_block:
        blocks.append({"lines": cur_block})

    for b in blocks:
        b_lines = b.get("lines") or []
        if not b_lines:
            continue

        left = min([float(ln["bbox"][0]) for ln in b_lines])
        right = max([float(ln["bbox"][2]) for ln in b_lines])
        top = min([float(ln["bbox"][1]) for ln in b_lines])
        bottom = max([float(ln["bbox"][3]) for ln in b_lines])
        b["bbox"] = [left, top, right, bottom]
        b["align"] = estimate_pdf_alignment(left, right, page_width)

        words_all: list[dict[str, Any]] = []
        for ln in b_lines:
            words_all.extend(list(ln.get("words") or []))
        b["font"] = summarize_font_from_words(words_all)

        fontnames = [w.get("fontname") for w in words_all]
        b["style"] = infer_style_from_fontnames(fontnames)

        b["text"] = "\n".join([str(ln.get("text", "")) for ln in b_lines]).strip()

    return blocks


# ==============================
# JSON schema helpers
# ==============================


def _ensure_common_schema(doc: dict[str, Any]) -> dict[str, Any]:
    """Ensure minimal common schema fields exist."""
    if "schema_version" not in doc:
        doc["schema_version"] = JSON_SCHEMA_VERSION
    if "pages" not in doc or not isinstance(doc["pages"], list):
        doc["pages"] = []
    return doc


def _wrap_page_texts(pages: list[str]) -> dict[str, Any]:
    return _ensure_common_schema(
        {
            "schema_version": JSON_SCHEMA_VERSION,
            "pages": [{"index": i + 1, "text": t} for i, t in enumerate(pages)],
        }
    )


def _prompt_for_password(path: str) -> Optional[str]:
    try:
        from .human_ask_tool import run_tool as human_ask
    except Exception:
        return None

    message = _(
        "prompt.password",
        default="Enter the password for this file:\n{path}",
    ).format(path=path)
    try:
        res_json = human_ask({"message": message, "is_password": True})
        res = json.loads(res_json)
    except Exception:
        return None

    pwd = str(res.get("user_reply") or "").strip()
    return pwd or None


def _looks_like_password_error(exc: Exception) -> bool:
    name = type(exc).__name__.lower()
    msg = str(exc).lower()
    return (
        "password" in name
        or "password" in msg
        or "encrypted" in msg
        or "decrypt" in msg
    )


def _load_pptx_pres(
    pptx_path: str, password: Optional[str] = None
) -> tuple[Optional[Any], list[str]]:
    warnings: list[str] = []
    if Presentation is None:
        warnings.append("python-pptx is not available")
        return None, warnings

    if msoffcrypto is not None:
        try:
            with open(pptx_path, "rb") as fin:
                office = msoffcrypto.OfficeFile(fin)
                encrypted = False
                try:
                    encrypted = bool(office.is_encrypted())
                except Exception:
                    encrypted = False

                if encrypted:
                    if not password:
                        password = _prompt_for_password(pptx_path)
                    if not password:
                        warnings.append("password is required for encrypted PPTX files")
                        return None, warnings

                    office.load_key(password=password)
                    buf = BytesIO()
                    office.decrypt(buf)
                    buf.seek(0)
                    return Presentation(buf), warnings
        except Exception as e:
            warnings.append(f"msoffcrypto failed: {type(e).__name__}: {e}")
            return None, warnings

    try:
        return Presentation(pptx_path), warnings
    except Exception as e:
        if password is None and _looks_like_password_error(e):
            prompted = _prompt_for_password(pptx_path)
            if prompted:
                return _load_pptx_pres(pptx_path, password=prompted)
        warnings.append(f"python-pptx failed: {type(e).__name__}: {e}")
        return None, warnings


# ==============================
# Extractors
# ==============================


def _extract_pdf_pages(
    pdf_path: str, password: str | None = None
) -> tuple[list[str], list[str]]:
    warnings: list[str] = []

    if pdfplumber is None:
        warnings.append("pdfplumber is not available")
        return [], warnings

    pages_text: list[str] = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                try:
                    words = page.extract_words(
                        use_text_flow=True,
                        keep_blank_chars=False,
                        extra_attrs=["fontname", "size"],
                    )
                except Exception:
                    words = []

                try:
                    blocks = build_pdf_blocks(words, page.width)
                    txt = "\n\n".join(
                        [b.get("text", "") for b in blocks if b.get("text")]
                    )
                    txt = normalize_text(txt)
                    if not txt.strip():
                        txt = normalize_text(page.extract_text() or "")
                except Exception:
                    txt = normalize_text(page.extract_text() or "")

                pages_text.append(txt)
    except Exception as e:
        warnings.append(f"pdfplumber failed: {type(e).__name__}: {e}")
        return [], warnings

    return pages_text, warnings


def _extract_pptx_pages(
    pptx_path: str, password: str | None = None
) -> tuple[list[str], list[str]]:
    warnings: list[str] = []

    if Presentation is None:
        warnings.append("python-pptx is not available")
        return [], warnings

    pages_text: list[str] = []
    try:
        pres = Presentation(pptx_path)
        for slide in pres.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                try:
                    if (
                        MSO_SHAPE_TYPE is not None
                        and shape.shape_type == MSO_SHAPE_TYPE.GROUP
                    ):
                        continue
                except Exception:
                    pass

                txt = ""
                try:
                    if (
                        getattr(shape, "has_text_frame", False)
                        and shape.text_frame is not None
                    ):
                        txt = shape.text_frame.text or ""
                except Exception:
                    txt = ""

                txt = normalize_text(txt)
                if txt.strip():
                    parts.append(txt)

            pages_text.append("\n".join(parts).strip())
    except Exception as e:
        warnings.append(f"python-pptx failed: {type(e).__name__}: {e}")
        return [], warnings

    return pages_text, warnings


# ==============================
# Main
# ==============================


def _read_common_json(path: str) -> tuple[dict[str, Any], list[str]]:
    warnings: list[str] = []
    try:
        data = json.loads(Path(path).read_text(encoding="utf-8"))
        if not isinstance(data, dict):
            return _ensure_common_schema({}), ["JSON root is not an object"]
        return _ensure_common_schema(data), warnings
    except Exception as e:
        return _ensure_common_schema({}), [
            f"Failed to read JSON: {type(e).__name__}: {e}"
        ]


def _get_pages_text(doc: dict[str, Any]) -> list[str]:
    pages = doc.get("pages")
    if not isinstance(pages, list):
        return []
    out: list[str] = []
    for p in pages:
        if not isinstance(p, dict):
            continue
        t = p.get("text")
        out.append(normalize_text(str(t or "")))
    return out


def run_tool(args: dict[str, Any]) -> str:
    path = (args.get("path") or "").strip()
    password = str(args.get("password") or "").strip() or None
    page_index = args.get("page_index") or args.get("page")
    max_chars = args.get("max_chars")

    if not path:
        return _("err.path_required", default="[read_pptx_pdf error] path is required")

    try:
        if max_chars is None:
            max_chars_i = DEFAULT_MAX_CHARS
        else:
            max_chars_i = int(max_chars)
            if max_chars_i <= 0:
                max_chars_i = DEFAULT_MAX_CHARS
    except Exception:
        max_chars_i = DEFAULT_MAX_CHARS

    p = Path(path)
    if not p.exists() or not p.is_file():
        return f"[read_pptx_pdf error] file not found: {path}"

    suffix = p.suffix.lower()
    warnings: list[str] = []

    if suffix == ".json":
        doc, warnings = _read_common_json(path)
        pages_text = _get_pages_text(doc)
    elif suffix == ".pdf":
        pages_text, warnings = _extract_pdf_pages(path, password=password)
        doc = _wrap_page_texts(pages_text)
    elif suffix == ".pptx":
        pages_text, warnings = _extract_pptx_pages(path, password=password)
        doc = _wrap_page_texts(pages_text)
    else:
        return f"[read_pptx_pdf error] unsupported file extension: {suffix}"

    pages_text = _get_pages_text(doc)

    if page_index is not None:
        try:
            idx = int(page_index)
        except Exception:
            return _(
                "err.page_index_int",
                default="[read_pptx_pdf error] page_index must be an integer",
            )

        if idx <= 0 or idx > len(pages_text):
            return f"[read_pptx_pdf error] page_index out of range: {idx}"

        out = pages_text[idx - 1]
    else:
        out = "\n\n".join([t for t in pages_text if t])

    if not out.strip():
        warnings.append("No text could be extracted")

    out = out[:max_chars_i]

    if warnings:
        warn_text = "\n".join([f"[read_pptx_pdf warn] {w}" for w in warnings])
        if out:
            return warn_text + "\n" + out
        return warn_text

    return out
