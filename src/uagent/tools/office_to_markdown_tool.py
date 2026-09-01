from __future__ import annotations

"""Convert common Office documents to readable Markdown."""

import json
import os
import shutil
import tempfile
from pathlib import Path
from typing import Any

from .i18n_helper import make_tool_translator

_ = make_tool_translator(__file__)

TOOL_SPEC: dict[str, Any] = {
    "tool_level": 1,
    "tool_genre": "office",
    "type": "function",
    "x_parallel_safe": True,
    "function": {
        "name": "office_to_markdown",
        "description": _(
            "tool.description",
            default="Convert PPTX, XLSX/XLSM, or DOCX files to Markdown. Excel number/date formats and basic emphasis are preserved; optionally add Mermaid flowcharts for recognizable process/dependency tables.",
        ),
        "x_search_terms": _(
            "x_search_terms",
            default=[
                "office to markdown",
                "pptx to markdown",
                "xlsx to markdown",
                "docx to markdown",
                "convert office file",
                "excel formatting",
                "mermaid flowchart",
            ],
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "input_path": {
                    "type": "string",
                    "description": _(
                        "param.input_path.description",
                        default="Input .pptx, .xlsx/.xlsm, or .docx path.",
                    ),
                },
                "password": {
                    "type": "string",
                    "description": _(
                        "param.password.description",
                        default="Password for an encrypted Office file (optional).",
                    ),
                },
                "output_path": {
                    "type": "string",
                    "description": _(
                        "param.output_path.description",
                        default="Optional output .md path.",
                    ),
                },
                "format": {
                    "type": "string",
                    "enum": ["auto", "pptx", "xlsx", "docx"],
                    "default": "auto",
                    "description": _(
                        "param.format.description",
                        default="Input format, normally auto-detected from the extension.",
                    ),
                },
                "include_notes": {
                    "type": "boolean",
                    "default": True,
                    "description": _(
                        "param.include_notes.description",
                        default="For PPTX, include speaker notes.",
                    ),
                },
                "overwrite": {
                    "type": "boolean",
                    "default": False,
                    "description": _(
                        "param.overwrite.description",
                        default="Allow replacing an existing output file.",
                    ),
                },
                "recalculate": {
                    "type": "boolean",
                    "default": True,
                    "description": _(
                        "param.recalculate.description",
                        default="For Excel files, recalculate formulas with Microsoft Excel before conversion when available.",
                    ),
                },
                "include_hidden": {
                    "type": "boolean",
                    "default": True,
                    "description": _(
                        "param.include_hidden.description",
                        default="For Excel files, include hidden rows and columns.",
                    ),
                },
                "include_mermaid": {
                    "type": "boolean",
                    "default": False,
                    "description": _(
                        "param.include_mermaid.description",
                        default="For Excel sheets with recognizable process or dependency columns, include a Mermaid flowchart.",
                    ),
                },
            },
            "required": ["input_path"],
            "additionalProperties": False,
        },
    },
}


def _cell(value: Any) -> str:
    return (
        ("" if value is None else str(value))
        .replace("\r", "")
        .replace("\n", "<br>")
        .replace("|", "\\|")
        .strip()
    )


def _table(rows: list[list[Any]]) -> str:
    rows = [[_cell(v) for v in row] for row in rows]
    rows = [row for row in rows if any(row)]
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    lines = [
        "| " + " | ".join(rows[0]) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in rows[1:])
    return "\n".join(lines)


def _resolve_encrypted(path: Path, password: str | None) -> tuple[Path, Path | None]:
    try:
        from .._pip_auto import install_with_status

        if not install_with_status("msoffcrypto-tool", "msoffcrypto"):
            return path, None
        import msoffcrypto
    except Exception as exc:
        raise RuntimeError(f"msoffcrypto could not be loaded: {exc}") from exc
    with path.open("rb") as source:
        office = msoffcrypto.OfficeFile(source)
        if not bool(office.is_encrypted()):
            return path, None
        if not password:
            raise ValueError("password is required for encrypted Office files")
        office.load_key(password=password)
        fd, name = tempfile.mkstemp(suffix=path.suffix)
        os.close(fd)
        temporary = Path(name)
        try:
            with temporary.open("wb") as target:
                office.decrypt(target)
        except Exception:
            temporary.unlink(missing_ok=True)
            raise
        return temporary, temporary


def _recalculate_xlsx(path: Path, password: str | None = None) -> dict[str, Any]:
    """Refresh Excel's cached formula values before extracting the workbook."""
    if os.name != "nt":
        return {
            "requested": True,
            "performed": False,
            "reason": "Microsoft Excel COM automation is available only on Windows",
        }

    try:
        import pythoncom  # type: ignore
        import win32com.client as win32  # type: ignore
    except Exception as exc:
        return {
            "requested": True,
            "performed": False,
            "reason": f"pywin32 or Microsoft Excel is unavailable: {exc}",
        }

    pythoncom.CoInitialize()
    excel = None
    workbook = None
    try:
        excel = win32.DispatchEx("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        try:
            excel.AskToUpdateLinks = False
        except Exception:
            pass
        if password:
            workbook = excel.Workbooks.Open(
                str(path), UpdateLinks=0, ReadOnly=False, Password=password
            )
        else:
            workbook = excel.Workbooks.Open(str(path), UpdateLinks=0, ReadOnly=False)
        excel.CalculateFullRebuild()
        try:
            excel.CalculateUntilAsyncQueriesDone()
        except Exception:
            pass
        workbook.Save()
        return {"requested": True, "performed": True}
    except Exception as exc:
        return {
            "requested": True,
            "performed": False,
            "reason": f"Excel recalculation failed: {type(exc).__name__}: {exc}",
        }
    finally:
        if workbook is not None:
            try:
                workbook.Close(SaveChanges=True)
            except Exception:
                pass
        if excel is not None:
            try:
                excel.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def _format_excel_value(value: Any, number_format: str) -> str:
    """Render common Excel number/date formats as users see them in Excel."""
    if value is None:
        return ""
    if hasattr(value, "strftime"):
        fmt = number_format or "yyyy-mm-dd"
        # Handle the most common Excel date tokens while retaining separators.
        result = fmt
        result = result.replace("yyyy", "%Y").replace("yy", "%y")
        result = result.replace("mm", "%m").replace("dd", "%d")
        result = result.replace("hh", "%H").replace("ss", "%S")
        result = result.replace("AM/PM", "%p")
        try:
            return value.strftime(result)
        except (TypeError, ValueError):
            return str(value)
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        fmt = number_format or "General"
        if fmt.lower() == "general":
            return str(value)
        percent = "%" in fmt
        number = value * 100 if percent else value
        # Count displayed decimal places, ignoring quoted literals.
        section = fmt.split(";")[0]
        decimals = (
            len(section.split(".", 1)[1].replace("%", "")) if "." in section else 0
        )
        grouping = "," in section.split(".", 1)[0]
        rendered = f"{number:{',' if grouping else ''}.{decimals}f}"
        if percent:
            rendered += "%"
        # Preserve common currency/accounting symbols from the format.
        for symbol in ("$", "€", "£", "¥", "￥"):
            if symbol in fmt:
                rendered = symbol + rendered
                break
        return rendered
    return str(value)


def _mermaid_flowchart(rows: list[list[Any]]) -> str:
    """Build a conservative flowchart from obvious source/target columns."""
    if len(rows) < 2 or not rows[0]:
        return ""
    headers = [str(value).strip().lower() for value in rows[0]]
    source_names = {"source", "from", "start", "parent", "親", "前工程", "前処理"}
    target_names = {"target", "to", "next", "child", "子", "次工程", "次処理", "依存先"}
    source_index = next((i for i, h in enumerate(headers) if h in source_names), None)
    target_index = next((i for i, h in enumerate(headers) if h in target_names), None)
    if source_index is None or target_index is None or source_index == target_index:
        return ""

    def clean(value: Any) -> str:
        text = str(value or "").strip().replace("\n", " ")
        return text.replace('"', "'").replace("[", "(").replace("]", ")")

    edges: list[tuple[str, str]] = []
    for row in rows[1:]:
        if max(source_index, target_index) >= len(row):
            continue
        source, target = clean(row[source_index]), clean(row[target_index])
        if source and target:
            edges.append((source, target))
    if not edges:
        return ""
    lines = ["```mermaid", "flowchart TD"]
    node_ids: dict[str, str] = {}
    for index, (source, target) in enumerate(edges):
        for label in (source, target):
            if label not in node_ids:
                node_ids[label] = f"N{len(node_ids) + 1}"
        lines.append(
            f'    {node_ids[source]}["{source}"] --> {node_ids[target]}["{target}"]'
        )
    lines.append("```")
    return "\n".join(lines)


def _convert_xlsx(
    path: Path,
    include_hidden: bool = True,
    title: str | None = None,
    include_mermaid: bool = False,
) -> str:
    try:
        from .exstruct_tool import _import_exstruct

        exstruct = _import_exstruct()
    except Exception as exc:
        raise RuntimeError(f"exstruct could not be loaded: {exc}") from exc
    if exstruct is None:
        raise RuntimeError("exstruct is required for Excel conversion")
    try:
        data = exstruct.extract(str(path), mode="standard", alpha_col=False)
        try:
            data = data.model_dump()
        except AttributeError:
            data = data.dict()
    except Exception as exc:
        raise RuntimeError(f"could not extract workbook with exstruct: {exc}") from exc
    parts = [f"# {title or path.stem}"]
    sheets = data.get("sheets", {}) if isinstance(data, dict) else {}
    # exstruct supplies values and formulas, while openpyxl supplies the
    # workbook's display formats (number format and basic emphasis styles).
    formatted_book = None
    try:
        from openpyxl import load_workbook

        formatted_book = load_workbook(str(path), data_only=True, read_only=True)
    except Exception:
        # Formatting is best-effort; conversion should still work with exstruct.
        formatted_book = None
    for name, sheet in sheets.items():
        rows = []
        worksheet = (
            formatted_book[name]
            if formatted_book is not None and name in formatted_book.sheetnames
            else None
        )
        for row in sheet.get("rows", []) if isinstance(sheet, dict) else []:
            cells = row.get("c", {}) if isinstance(row, dict) else {}
            if isinstance(cells, dict):
                keys = sorted(
                    cells,
                    key=lambda value: (
                        int(value) if str(value).isdigit() else str(value)
                    ),
                )
                row_number = int(row.get("r", 0) or 0) if isinstance(row, dict) else 0
                if (
                    worksheet is not None
                    and not include_hidden
                    and row_number
                    and worksheet.row_dimensions[row_number].hidden
                ):
                    continue
                rendered = []
                for key in keys:
                    value = cells[key]
                    column_number = int(key) + 1 if str(key).isdigit() else 0
                    if (
                        worksheet is not None
                        and not include_hidden
                        and column_number
                        and worksheet.cell(row_number, column_number).column_letter
                        in worksheet.column_dimensions
                        and worksheet.column_dimensions[
                            worksheet.cell(row_number, column_number).column_letter
                        ].hidden
                    ):
                        continue
                    if worksheet is not None and row_number and column_number:
                        excel_cell = worksheet.cell(row_number, column_number)
                        if excel_cell.value is not None:
                            value = _format_excel_value(
                                excel_cell.value, excel_cell.number_format
                            )
                            if excel_cell.font.bold:
                                value = f"**{value}**"
                            if excel_cell.font.italic:
                                value = f"*{value}*"
                    rendered.append(value)
                rows.append(rendered)
        table = _table(rows)
        parts.append(f"## {name}\n\n{table or '_（データなし）_'}")
        if include_mermaid:
            flowchart = _mermaid_flowchart(rows)
            if flowchart:
                parts.append(f"### {name} flowchart\n\n{flowchart}")
        if isinstance(sheet, dict) and sheet.get("merged_ranges"):
            parts.append(
                f"Merged ranges: {', '.join(map(str, sheet['merged_ranges']))}"
            )
        if isinstance(sheet, dict) and sheet.get("formulas_map"):
            parts.append(f"Formulas: {len(sheet['formulas_map'])}")
    return "\n\n".join(parts).rstrip() + "\n"


def _convert_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required for Word conversion") from exc
    doc = Document(str(path))
    parts = [f"# {path.stem}"]
    tables = {table._element: table for table in doc.tables}
    paragraphs = {para._p: para for para in doc.paragraphs}
    counters: dict[str, int] = {}
    table_index = 0
    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    for element in doc.element.body.iterchildren():
        if element in tables:
            table_index += 1
            rendered = _table(
                [[cell.text for cell in row.cells] for row in tables[element].rows]
            )
            if rendered:
                parts.append(f"### Table {table_index}\n\n{rendered}")
            continue
        para = paragraphs.get(element)
        if para is None or not para.text.strip():
            continue
        text = para.text.strip()
        style_obj = getattr(para, "style", None)
        style_id = str(getattr(style_obj, "style_id", "") or "").lower()
        style = str(getattr(style_obj, "name", "") or "").lower()
        level = None
        for prefix in ("heading", "見出し"):
            if style_id.startswith(prefix) or style.startswith(prefix + " "):
                try:
                    level = int(
                        (style_id if style_id.startswith(prefix) else style)
                        .split()[-1]
                        .replace(prefix, "")
                    )
                except ValueError:
                    level = 2
                break
        if level is not None:
            parts.append(f"{'#' * min(max(level, 1), 6)} {text}")
        elif "list bullet" in style or "箇条書き" in style:
            parts.append(f"- {text}")
        else:
            num_id = para._p.find(f"{ns}numPr/{ns}numId")
            if num_id is not None:
                key = num_id.get(f"{ns}val", "default")
                counters[key] = counters.get(key, 0) + 1
                parts.append(f"{counters[key]}. {text}")
            else:
                parts.append(text)
    return "\n\n".join(parts).rstrip() + "\n"


def _convert_pptx(path: Path, include_notes: bool = True) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PowerPoint conversion") from exc
    presentation = Presentation(str(path))
    parts = [f"# {path.stem}"]

    def visit(shape: Any, texts: list[str]) -> None:
        if getattr(shape, "has_table", False):
            rendered = _table(
                [[cell.text for cell in row.cells] for row in shape.table.rows]
            )
            if rendered:
                texts.append(rendered)
        elif getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                texts.append(text)
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                visit(child, texts)

    for number, slide in enumerate(presentation.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            visit(shape, texts)
        block = [f"## Slide {number}"] + texts
        if include_notes:
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes = ""
            if notes:
                block.extend(["### Notes", notes])
        parts.append("\n\n".join(block))
    return "\n\n".join(parts).rstrip() + "\n"


def run_tool(args: dict[str, Any]) -> str:
    args = args or {}
    input_path = str(args.get("input_path") or "").strip()
    if not input_path:
        return json.dumps(
            {
                "ok": False,
                "error": _("err.input_path_required", default="input_path is required"),
            },
            ensure_ascii=False,
        )
    path = Path(input_path)
    if not path.is_file():
        return json.dumps(
            {
                "ok": False,
                "error": _(
                    "err.file_not_found", default="file not found: {path}"
                ).format(path=input_path),
            },
            ensure_ascii=False,
        )
    detected = {".pptx": "pptx", ".xlsx": "xlsx", ".xlsm": "xlsx", ".docx": "docx"}.get(
        path.suffix.lower()
    )
    fmt = str(args.get("format") or "auto").lower()
    if fmt == "auto":
        fmt = detected or ""
    if detected != fmt:
        return json.dumps(
            {
                "ok": False,
                "error": _(
                    "err.format_mismatch",
                    default="unsupported or mismatched format; use .pptx, .xlsx/.xlsm, or .docx",
                ),
            },
            ensure_ascii=False,
        )
    temporary: Path | None = None
    recalc_temporary: Path | None = None
    recalculation: dict[str, Any] | None = None
    try:
        password = str(args.get("password") or "").strip() or None
        readable_path, temporary = _resolve_encrypted(path, password)
        if fmt == "pptx":
            markdown = _convert_pptx(
                readable_path, args.get("include_notes", True) is not False
            )
        elif fmt == "xlsx":
            conversion_path = readable_path
            if args.get("recalculate", True) is not False:
                fd, recalc_name = tempfile.mkstemp(suffix=readable_path.suffix)
                os.close(fd)
                recalc_temporary = Path(recalc_name)
                shutil.copy2(readable_path, recalc_temporary)
                recalculation = _recalculate_xlsx(
                    recalc_temporary, None if temporary is not None else password
                )
                conversion_path = recalc_temporary
            markdown = _convert_xlsx(
                conversion_path,
                args.get("include_hidden", True) is not False,
                title=path.stem,
                include_mermaid=args.get("include_mermaid", False) is True,
            )
        else:
            markdown = _convert_docx(readable_path)
        output_path = str(args.get("output_path") or "").strip()
        if output_path:
            out = Path(output_path)
            if out.resolve() == path.resolve():
                raise ValueError(
                    _(
                        "err.input_output_same",
                        default="output_path must not be the input file",
                    )
                )
            if out.exists() and not args.get("overwrite", False):
                raise FileExistsError(
                    _(
                        "err.output_exists",
                        default="output file already exists: {path}; set overwrite=true to replace it",
                    ).format(path=output_path)
                )
            out.parent.mkdir(parents=True, exist_ok=True)
            out.write_text(markdown, encoding="utf-8")
        result: dict[str, Any] = {
            "ok": True,
            "format": fmt,
            "input_path": str(path),
            "output_path": output_path or None,
            "markdown": markdown,
        }
        if recalculation is not None:
            result["recalculation"] = recalculation
        return json.dumps(result, ensure_ascii=False)
    except Exception as exc:
        return json.dumps({"ok": False, "error": str(exc)}, ensure_ascii=False)
    finally:
        for temporary_path in (recalc_temporary, temporary):
            if temporary_path is not None:
                try:
                    temporary_path.unlink(missing_ok=True)
                except OSError:
                    # A third-party extractor may briefly retain a file handle.
                    pass
